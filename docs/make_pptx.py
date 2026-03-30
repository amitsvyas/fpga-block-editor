"""
Generate fpga_block_editor_overview.pptx
Run: python3 make_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette (matches the app's dark theme) ──────────────────────────
C_BG        = RGBColor(0x1e, 0x1e, 0x2e)   # dark navy
C_CARD      = RGBColor(0x25, 0x25, 0x38)   # node card colour
C_INDIGO    = RGBColor(0x4f, 0x46, 0xe5)   # primary accent
C_INDIGO_LT = RGBColor(0x81, 0x81, 0xf8)   # lighter indigo for headings
C_WHITE     = RGBColor(0xe2, 0xe8, 0xf0)
C_MUTED     = RGBColor(0x71, 0x71, 0x7a)
C_YELLOW    = RGBColor(0xfa, 0xcc, 0x15)   # clock signal
C_PURPLE    = RGBColor(0xa7, 0x8b, 0xfa)   # AXI signal
C_ORANGE    = RGBColor(0xfb, 0x92, 0x3c)   # control signal
C_GREEN     = RGBColor(0x4a, 0xde, 0x80)
C_RED       = RGBColor(0xf8, 0x71, 0x71)

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]   # completely blank


def add_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=C_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill=None, border=None, border_w=Pt(1)):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = border_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def accent_bar(slide, color=C_INDIGO):
    box(slide, 0, 0, Inches(0.08), H, fill=color, border=color)


def slide_label(slide, label, color=C_INDIGO_LT):
    txt(slide, label.upper(),
        Inches(0.2), Inches(0.18), Inches(5), Inches(0.35),
        size=9, bold=True, color=color)


def divider(slide, y, color=C_CARD):
    box(slide, Inches(0.2), y, W - Inches(0.4), Inches(0.02), fill=color)


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)

# Big indigo circle decoration
c = s.shapes.add_shape(9, Inches(8.5), Inches(-1.2), Inches(5.5), Inches(5.5))  # oval
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x2d, 0x2d, 0x55)
c.line.fill.background()

txt(s, "LATTICE AVANT", Inches(0.5), Inches(1.6), Inches(9), Inches(0.8),
    size=13, bold=True, color=C_INDIGO_LT)
txt(s, "FPGA Block Editor",
    Inches(0.5), Inches(2.15), Inches(10), Inches(1.4),
    size=48, bold=True, color=C_WHITE)
txt(s, "Visual system-level design tool for Field & Regional Application Engineers",
    Inches(0.5), Inches(3.65), Inches(8.5), Inches(0.7),
    size=18, color=RGBColor(0xa1, 0xa1, 0xaa))
divider(s, Inches(4.55))
txt(s, "Web-based  ·  No install  ·  Exports Verilog, VHDL & JSON netlist",
    Inches(0.5), Inches(4.65), Inches(9), Inches(0.5),
    size=14, color=C_MUTED)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
slide_label(s, "The Problem")

txt(s, "Today's workflow has gaps",
    Inches(0.3), Inches(0.55), Inches(9), Inches(0.7),
    size=32, bold=True, color=C_WHITE)

problems = [
    ("📊", "PowerPoint diagrams",    "Blocks and arrows with no defined port types, bus widths, or clock domains"),
    ("🔌", "No connectivity rules",  "Nothing stops you drawing clock→data or output→output connections"),
    ("📤", "Manual translation",     "Engineers re-enter the same topology by hand in Radiant or Vivado"),
    ("🔄", "Version drift",          "The diagram and the actual HDL quickly go out of sync"),
]

for i, (icon, title, desc) in enumerate(problems):
    col = i % 2
    row = i // 2
    lx  = Inches(0.3  + col * 6.4)
    ty  = Inches(1.55 + row * 2.3)
    box(s, lx, ty, Inches(6.1), Inches(2.0), fill=C_CARD,
        border=RGBColor(0x3a, 0x3a, 0x5c))
    txt(s, icon,  lx + Inches(0.18), ty + Inches(0.18), Inches(0.6), Inches(0.6), size=22)
    txt(s, title, lx + Inches(0.7),  ty + Inches(0.22), Inches(5.0), Inches(0.5),
        size=16, bold=True, color=C_WHITE)
    txt(s, desc,  lx + Inches(0.18), ty + Inches(0.82), Inches(5.7), Inches(1.0),
        size=12, color=RGBColor(0xa1, 0xa1, 0xaa))


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
slide_label(s, "The Solution")

txt(s, "One tool — from diagram to RTL",
    Inches(0.3), Inches(0.55), Inches(10), Inches(0.7),
    size=32, bold=True, color=C_WHITE)

steps = [
    (C_INDIGO,    "1", "Drag blocks",     "Choose from 70+ component types across Lattice primitives, OSS IPs, 3rd-party chips, and power regulators"),
    (C_PURPLE,    "2", "Draw connections","Ports are typed and directional — invalid wires are rejected automatically"),
    (C_GREEN,     "3", "Configure",       "Click any block to set instance name, clock domain, Fmax target, and edit/add ports"),
    (C_YELLOW,    "4", "Export",          "One click: JSON netlist, structural Verilog (.v), or structural VHDL (.vhd) — ready for Radiant or Vivado"),
]

for i, (color, num, title, desc) in enumerate(steps):
    lx = Inches(0.25 + i * 3.2)
    ty = Inches(1.5)
    # Arrow connector (except last)
    if i < 3:
        box(s, lx + Inches(2.85), ty + Inches(0.7), Inches(0.35), Inches(0.08),
            fill=C_MUTED)
    box(s, lx, ty, Inches(2.8), Inches(4.8), fill=C_CARD,
        border=color, border_w=Pt(1.5))
    # Number badge
    badge = s.shapes.add_shape(9, lx + Inches(1.05), ty + Inches(0.2),
                                Inches(0.7), Inches(0.7))
    badge.fill.solid(); badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    txt(s, num, lx + Inches(1.05), ty + Inches(0.18), Inches(0.7), Inches(0.7),
        size=18, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, title, lx + Inches(0.15), ty + Inches(1.1), Inches(2.5), Inches(0.55),
        size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, desc,  lx + Inches(0.15), ty + Inches(1.75), Inches(2.5), Inches(2.8),
        size=11, color=RGBColor(0xa1, 0xa1, 0xaa), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — Component Library
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
slide_label(s, "Component Library")

txt(s, "70+ components across 6 categories",
    Inches(0.3), Inches(0.55), Inches(10), Inches(0.7),
    size=32, bold=True, color=C_WHITE)

categories = [
    (RGBColor(0x6366,0xf1,0x00)[:3] if False else RGBColor(0x63, 0x66, 0xf1),
     "Primitives",       "Generic IP · PLL/Clock · I/O Pad\nDSP Block · RAM Block",              5),
    (RGBColor(0x0e, 0xa5, 0xe9),
     "PCIe & High-Speed","SerDes PMA · PCIe PCS · PCIe Ctrl\nEthernet MAC · MIPI D-PHY",         5),
    (RGBColor(0x10, 0xb9, 0x81),
     "Memory Controllers","DDR3 · DDR4 · LPDDR4",                                                3),
    (RGBColor(0x8b, 0x5c, 0xf6),
     "Open Source IPs",  "VexRiscv · PicoRV32 · LiteDRAM\nLiteEth · AXI Crossbar · AES · FFT", 16),
    (RGBColor(0xf5, 0x9e, 0x0b),
     "Power",            "LDO · Buck · Boost · Buck-Boost\nMulti-Phase · PMBus · Sequencer",     7),
    (RGBColor(0xec, 0x48, 0x99),
     "3rd Party Chips",  "Ethernet PHY · DDR3 SDRAM · HyperRAM\nMIPI Camera · Radar · IMU · HDMI",26),
]

cols = 3
for i, (color, cat, items, count) in enumerate(categories):
    col = i % cols
    row = i // cols
    lx  = Inches(0.25 + col * 4.33)
    ty  = Inches(1.5  + row * 2.6)
    box(s, lx, ty, Inches(4.1), Inches(2.35), fill=C_CARD,
        border=color, border_w=Pt(1.5))
    # Colour strip top
    box(s, lx, ty, Inches(4.1), Inches(0.08), fill=color)
    txt(s, cat, lx + Inches(0.15), ty + Inches(0.15), Inches(3.0), Inches(0.45),
        size=14, bold=True, color=color)
    # Count badge
    bdg = s.shapes.add_shape(1, lx + Inches(3.35), ty + Inches(0.2),
                               Inches(0.6), Inches(0.28))
    bdg.fill.solid(); bdg.fill.fore_color.rgb = color
    bdg.line.fill.background()
    txt(s, f"{count} blocks", lx + Inches(3.35), ty + Inches(0.2),
        Inches(0.6), Inches(0.28), size=8, bold=True,
        color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, items, lx + Inches(0.15), ty + Inches(0.65), Inches(3.8), Inches(1.5),
        size=11, color=RGBColor(0xa1, 0xa1, 0xaa))


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — Canvas Features
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
slide_label(s, "Canvas Features")

txt(s, "Smart connections, not just lines",
    Inches(0.3), Inches(0.55), Inches(10), Inches(0.7),
    size=32, bold=True, color=C_WHITE)

# Left column — feature list
features = [
    (C_YELLOW,  "Typed ports",          "Every port has a direction (in/out/bidir) and signal type (clock, data, control, AXI, power)"),
    (C_RED,     "Connection validation","Output→output and clock→data connections are blocked at draw time"),
    (C_PURPLE,  "Bus width encoding",   "Wire thickness scales with bit-width — 1-bit vs 256-bit buses are visually distinct"),
    (C_GREEN,   "Snap to grid",         "16px grid keeps diagrams tidy without manual alignment"),
    (C_ORANGE,  "MiniMap",              "Thumbnail overview for navigating large multi-block designs"),
]
for i, (color, title, desc) in enumerate(features):
    ty = Inches(1.55 + i * 1.05)
    box(s, Inches(0.25), ty + Inches(0.12), Inches(0.05), Inches(0.55), fill=color)
    txt(s, title, Inches(0.45), ty + Inches(0.1), Inches(3.8), Inches(0.38),
        size=13, bold=True, color=color)
    txt(s, desc,  Inches(0.45), ty + Inches(0.48), Inches(5.8), Inches(0.5),
        size=11, color=RGBColor(0xa1, 0xa1, 0xaa))

# Right column — signal legend
box(s, Inches(7.0), Inches(1.5), Inches(5.9), Inches(5.5), fill=C_CARD,
    border=RGBColor(0x3a, 0x3a, 0x5c))
txt(s, "Signal Type Legend", Inches(7.2), Inches(1.65), Inches(5.0), Inches(0.45),
    size=13, bold=True, color=C_WHITE)
divider(s, Inches(2.2))

signals = [
    (C_YELLOW, "Clock",   "Dashed wire · triggers sequential logic"),
    (C_WHITE,  "Data",    "Width-encoded · general purpose buses"),
    (C_ORANGE, "Control", "Enable, reset, chip-select signals"),
    (C_PURPLE, "AXI",     "AXI4 / AXI4-Lite / AXI4-Stream buses"),
    (C_RED,    "Power",   "VDD, GND, power rails"),
]
for i, (color, name, desc) in enumerate(signals):
    ty = Inches(2.35 + i * 0.85)
    circ = s.shapes.add_shape(9, Inches(7.2), ty + Inches(0.12),
                               Inches(0.28), Inches(0.28))
    circ.fill.solid(); circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    txt(s, name, Inches(7.6),  ty + Inches(0.06), Inches(1.2), Inches(0.35),
        size=13, bold=True, color=color)
    txt(s, desc, Inches(7.6),  ty + Inches(0.38), Inches(5.0), Inches(0.38),
        size=10, color=C_MUTED)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — Export Formats
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
slide_label(s, "Export Formats")

txt(s, "From diagram to RTL in one click",
    Inches(0.3), Inches(0.55), Inches(10), Inches(0.7),
    size=32, bold=True, color=C_WHITE)

formats = [
    (C_GREEN,   "JSON Netlist",   ".json",
     "Structured machine-readable netlist.\n\nContains all modules (ip_type, instance_name, ports, clock_domain, fmax) and all nets (source/target instance+port, bus width, signal type).\n\nUse for: version control, design reviews, custom automation scripts."),
    (C_INDIGO_LT, "Verilog",     ".v",
     "Structural Verilog top-level.\n\n• Black-box module declarations\n• Wire declarations (width-correct)\n• Named port-map instantiations\n• Inline vendor primitive hints\n  (Radiant ↔ Vivado equivalents)\n\nImport directly into Lattice Radiant or Xilinx Vivado as a design source."),
    (C_PURPLE,  "VHDL",          ".vhd",
     "Structural VHDL — entity / architecture structural.\n\n• Component declarations per module type\n• Signal declarations\n• Port-map instantiations\n• Unconnected ports → open\n\nImport directly into Lattice Radiant or Xilinx Vivado as a design source."),
]

for i, (color, fmt, ext, desc) in enumerate(formats):
    lx = Inches(0.25 + i * 4.33)
    ty = Inches(1.45)
    box(s, lx, ty, Inches(4.1), Inches(5.6), fill=C_CARD,
        border=color, border_w=Pt(1.5))
    box(s, lx, ty, Inches(4.1), Inches(0.08), fill=color)
    # Extension badge
    bdg = s.shapes.add_shape(1, lx + Inches(0.15), ty + Inches(0.2),
                               Inches(0.65), Inches(0.3))
    bdg.fill.solid(); bdg.fill.fore_color.rgb = color
    bdg.line.fill.background()
    txt(s, ext, lx + Inches(0.15), ty + Inches(0.2), Inches(0.65), Inches(0.3),
        size=10, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, fmt, lx + Inches(0.95), ty + Inches(0.2), Inches(2.9), Inches(0.42),
        size=16, bold=True, color=color)
    txt(s, desc, lx + Inches(0.18), ty + Inches(0.78), Inches(3.75), Inches(4.6),
        size=11, color=RGBColor(0xa1, 0xa1, 0xaa))


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — AI Agent Potential
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
slide_label(s, "Future: AI Agent")

txt(s, "Natural language → block diagram",
    Inches(0.3), Inches(0.55), Inches(10), Inches(0.7),
    size=32, bold=True, color=C_WHITE)

# Prompt box (left)
box(s, Inches(0.3), Inches(1.5), Inches(5.5), Inches(1.4),
    fill=RGBColor(0x1a, 0x1a, 0x3e), border=C_INDIGO)
txt(s, "Engineer types:",
    Inches(0.5), Inches(1.6), Inches(5), Inches(0.35),
    size=10, color=C_MUTED)
txt(s, '"Design a PCIe Gen3 system with DDR4,\nGigabit Ethernet, and a RISC-V processor"',
    Inches(0.5), Inches(1.95), Inches(5.2), Inches(0.85),
    size=13, italic=True, color=C_INDIGO_LT)

# Arrow
txt(s, "→", Inches(5.95), Inches(2.0), Inches(0.8), Inches(0.6),
    size=28, color=C_INDIGO, align=PP_ALIGN.CENTER)

# Result box (right)
box(s, Inches(6.8), Inches(1.5), Inches(6.1), Inches(1.4),
    fill=RGBColor(0x1a, 0x1a, 0x3e), border=C_GREEN)
txt(s, "AI agent auto-populates:",
    Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.35),
    size=10, color=C_MUTED)
txt(s, "PCIe Ctrl · SerDes PMA · DDR4 Ctrl\nDDR4 SDRAM · VexRiscv · Ethernet MAC\nEthernet PHY — all wired correctly",
    Inches(7.0), Inches(1.95), Inches(5.8), Inches(0.85),
    size=12, color=C_GREEN)

# How it would work
txt(s, "How it would work",
    Inches(0.3), Inches(3.15), Inches(6), Inches(0.45),
    size=15, bold=True, color=C_WHITE)

steps = [
    (C_INDIGO,  "Claude API", "Engineer describes the design in plain English via a chat panel in the tool"),
    (C_PURPLE,  "Tool use",   "Claude calls structured tools to addNode(), addEdge(), setPort() — same API as the canvas"),
    (C_GREEN,   "Validation", "Connection rules still enforced — AI cannot produce invalid wiring"),
    (C_YELLOW,  "Export",     "Engineer reviews the generated diagram and exports Verilog/VHDL as normal"),
]
for i, (color, title, desc) in enumerate(steps):
    lx = Inches(0.3 + i * 3.25)
    ty = Inches(3.65)
    box(s, lx, ty, Inches(3.05), Inches(3.4), fill=C_CARD,
        border=color, border_w=Pt(1.2))
    txt(s, title, lx + Inches(0.15), ty + Inches(0.2), Inches(2.75), Inches(0.45),
        size=13, bold=True, color=color)
    txt(s, desc,  lx + Inches(0.15), ty + Inches(0.72), Inches(2.75), Inches(2.5),
        size=11, color=RGBColor(0xa1, 0xa1, 0xaa))

txt(s, "⚡  Backend needed: Next.js API route + Anthropic SDK (~1 day of work)",
    Inches(0.3), Inches(7.05), Inches(12), Inches(0.35),
    size=11, color=C_MUTED, italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — Deployment
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
slide_label(s, "Deployment Options")

txt(s, "Zero infrastructure required",
    Inches(0.3), Inches(0.55), Inches(10), Inches(0.7),
    size=32, bold=True, color=C_WHITE)

options = [
    (C_INDIGO_LT, "Vercel (live now)",
     "Free tier · GitHub-connected\nAuto-redeploys on every commit\nAccess from any device, anywhere\n\nURL: fpga-block-editor.vercel.app",
     "Recommended for sharing with\nremote / field AEs"),
    (C_GREEN, "Static file share\n(internal network)",
     "Copy out/ folder to a Windows\nfile share or IIS web root\nNo server software needed\nWorks fully offline\n\nOpen index.html in Edge",
     "Best for air-gapped\nor on-premise environments"),
    (C_ORANGE, "Windows laptop\n(local dev)",
     "Install Node.js 22 from nodejs.org\nnpm install\nnpm run dev\nOpens at localhost:3000\n\nFull live-reload for customisation",
     "Best for engineers who want\nto modify or extend the tool"),
    (C_PURPLE, "Azure Static\nWeb Apps",
     "Enterprise Azure tenant\nSupports AAD / SSO integration\nBuilt-in CI/CD from GitHub\nFree tier available\n\naz staticwebapp create ...",
     "Best for enterprise with\nexisting Azure footprint"),
]

for i, (color, title, body, note) in enumerate(options):
    col = i % 2
    row = i // 2
    lx  = Inches(0.25 + col * 6.5)
    ty  = Inches(1.45 + row * 2.85)
    box(s, lx, ty, Inches(6.2), Inches(2.6), fill=C_CARD,
        border=color, border_w=Pt(1.5))
    box(s, lx, ty, Inches(6.2), Inches(0.07), fill=color)
    txt(s, title, lx + Inches(0.18), ty + Inches(0.15), Inches(4.5), Inches(0.6),
        size=14, bold=True, color=color)
    txt(s, body,  lx + Inches(0.18), ty + Inches(0.72), Inches(3.7), Inches(1.7),
        size=11, color=RGBColor(0xa1, 0xa1, 0xaa))
    # Note pill (right side)
    box(s, lx + Inches(3.9), ty + Inches(0.72), Inches(2.1), Inches(1.5),
        fill=RGBColor(0x12, 0x12, 0x1f), border=color, border_w=Pt(0.75))
    txt(s, note, lx + Inches(4.0), ty + Inches(0.82), Inches(1.9), Inches(1.3),
        size=10, color=color, italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — Summary / Next Steps
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)

# Large background decoration
c2 = s.shapes.add_shape(9, Inches(9.5), Inches(3.5), Inches(6), Inches(6))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x38)
c2.line.fill.background()

txt(s, "Summary", Inches(0.3), Inches(0.55), Inches(5), Inches(0.55),
    size=32, bold=True, color=C_WHITE)

bullets = [
    (C_GREEN,      "✓  Live on Vercel",           "Shareable URL — no install for AEs"),
    (C_GREEN,      "✓  Offline zip available",    "Works from a Windows file share or USB drive"),
    (C_GREEN,      "✓  70+ components",           "Lattice primitives, OSS IPs, power, 3rd-party chips"),
    (C_GREEN,      "✓  Verilog + VHDL export",    "Structural RTL ready for Radiant or Vivado"),
    (C_INDIGO_LT,  "◎  AI agent (next step)",     "Claude API → natural language → auto-generated diagrams"),
    (C_ORANGE,     "◎  Import / load design",     "Load a previously saved JSON back onto the canvas"),
    (C_PURPLE,     "◎  Constraint export",        ".xdc / .pdc timing constraint file from clock domains + Fmax"),
]

for i, (color, title, desc) in enumerate(bullets):
    ty = Inches(1.35 + i * 0.77)
    box(s, Inches(0.25), ty + Inches(0.15), Inches(0.06), Inches(0.45), fill=color)
    txt(s, title, Inches(0.42), ty + Inches(0.1), Inches(4.2), Inches(0.4),
        size=13, bold=True, color=color)
    txt(s, desc,  Inches(0.42), ty + Inches(0.45), Inches(6.8), Inches(0.35),
        size=11, color=C_MUTED)

txt(s, "Built with Next.js · React Flow · Zustand · TypeScript",
    Inches(0.25), Inches(7.1), Inches(9), Inches(0.35),
    size=10, color=RGBColor(0x3f, 0x3f, 0x5c), italic=True)


# ── Save ──────────────────────────────────────────────────────────
out_path = "fpga_block_editor_overview.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
